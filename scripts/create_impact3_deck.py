"""
Create Impact3 AI Audit Presentation Deck
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RgbColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RgbColor(0x41, 0x69, 0xe1)
WHITE = RgbColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RgbColor(0xf5, 0xf5, 0xf5)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if point.startswith("##"):
            p.text = point.replace("## ", "")
            p.font.size = Pt(22)
            p.font.bold = True
            p.space_before = Pt(16)
        elif point.startswith("#"):
            p.text = point.replace("# ", "")
            p.font.size = Pt(26)
            p.font.bold = True
        else:
            p.text = "• " + point if not point.startswith("-") else point.replace("-", "  •")
            p.font.size = Pt(18)
            p.space_before = Pt(8)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.9), Inches(5.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if point.startswith("##"):
            p.text = point.replace("## ", "")
            p.font.bold = True
            p.font.size = Pt(20)
        else:
            p.text = "• " + point
            p.font.size = Pt(16)
        p.space_before = Pt(6)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.2), Inches(5.9), Inches(5.8))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if point.startswith("##"):
            p.text = point.replace("## ", "")
            p.font.bold = True
            p.font.size = Pt(20)
        else:
            p.text = "• " + point
            p.font.size = Pt(16)
        p.space_before = Pt(6)

    return slide

# SLIDE 1: Title
add_title_slide(prs, "Impact3 AI & Automation Audit", "Discovery Findings & Strategic Recommendations\n\nPrepared by Reprise AI | January 2026")

# SLIDE 2: Agenda
add_content_slide(prs, "Agenda", [
    "Executive Summary",
    "Your Team's Voice - Interview Insights",
    "Pain Points Matrix",
    "Automation Opportunities",
    "Recommended Solutions",
    "Implementation Roadmap",
    "ROI Projections",
    "Next Steps"
])

# SLIDE 3: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "## Goal",
    "Scale from $3M to $6M ARR without doubling headcount (40 → 60 people)",
    "",
    "## Interviews Completed: 10/10 departments",
    "",
    "## Key Finding",
    "Communication overhead and manual data workflows are the primary scaling blockers",
    "",
    "## Opportunity",
    "70%+ efficiency gains possible in critical areas",
    "",
    "## Recommendation",
    "Phased automation rollout starting with Discord intelligence and report automation"
])

# SLIDE 4: Company at a Glance
add_content_slide(prs, "Your Company at a Glance", [
    "## Current State",
    "ARR: $3M | Team Size: 40 | Manual Hours/Week: 20-40",
    "",
    "## Target State",
    "ARR: $6M | Team Size: 60 (with efficiency) | Manual Hours: 50% reduction",
    "",
    "## Services",
    "Social Media, Content, PR, Paid Growth for Web3/Crypto clients",
    "",
    "## Tools in Use",
    "Notion, Discord, n8n, Typefully, Figma, Toggle",
    "",
    "## Existing Automation",
    "Zapier + n8n (scoreboards, analytics → Google Sheets)"
])

# SLIDE 5: Interview Coverage
add_content_slide(prs, "Interview Coverage - 10 Departments", [
    "Raul (CEO) - Scaling, Discord overload",
    "Jeff (CSO) - Multi-platform coordination",
    "Matt (Social Media) - Content discovery, research",
    "Leah (Copywriting) - Asset findability",
    "Jess (Visual Production) - Incomplete briefs",
    "Joe (Paid Growth) - Competitor analysis",
    "Kyle R (Sales/Co-Owner) - Pipeline automation",
    "Kyle H + Patrick (PR) - Placement tracking",
    "Diana + team (Account Management) - Report compilation",
    "Gonzalo (KOL) - Influencer coordination"
])

# SLIDE 6: Critical Pain Points
add_title_slide(prs, "Pain Points Matrix", "Critical Priority Issues")

# SLIDE 7: Discord Overload
add_content_slide(prs, "CRITICAL: Discord Communication Overload", [
    "## Who: CEO + entire team",
    "",
    "## Scale",
    "50+ messages/day for CEO alone",
    "20 messages morning + 30-40 throughout day",
    "",
    "## Impact",
    '"Shit ton of hours" - CEO quote',
    "",
    "## Scaling Risk",
    '"Adding 20-40 people, I don\'t know how anyone will keep up with Discord"',
    "",
    "## Target: 70% reduction in back-and-forth"
])

# SLIDE 8: Research Inefficiency
add_content_slide(prs, "CRITICAL: Research Inefficiency", [
    "## Who: Social Media, PR, Strategy",
    "",
    "## Scale",
    "10-20 hrs/week on manual research (PR alone)",
    "",
    "## Current Tool",
    "Internal AI research tool 'not fully working'",
    "",
    "## Goal",
    "10-minute morning scroll = all client research done",
    "",
    "## CEO Quote",
    '"What would you automate across Impact3? Research."'
])

# SLIDE 9: High Priority Pain Points
add_two_column_slide(prs, "HIGH PRIORITY Pain Points",
    [
        "## Incomplete Creative Briefs",
        "Who: Visual Production (Jess)",
        "Impact: 5 min → 2+ HOURS delay",
        "Root Cause: Missing references, platforms, copy, dimensions",
        "Solution: Required field validation"
    ],
    [
        "## Manual Report Compilation",
        "Who: Account Managers (Diana)",
        "Time Cost: 2-4 hrs/week",
        "Connection: CEO's primary visibility",
        'Raul: "If we can automate that, that\'s even better"'
    ]
)

# SLIDE 10: Medium Priority
add_two_column_slide(prs, "MEDIUM PRIORITY Pain Points",
    [
        "## Reply Gang Content Discovery",
        "Who: Social Media team",
        "Challenge: Finding relevant tweets",
        "Current Tool: X Radar (no API)",
        'Matt: "If you can scrape X Radar, that would be crazy"'
    ],
    [
        "## Cross-Platform Fragmentation",
        "Who: All client-facing teams",
        "Issue: Slack + Telegram + Teams + Discord",
        "Impact: Context switching, missed messages",
        "Solution: Unified inbox or routing"
    ]
)

# SLIDE 11: Cross-Department Patterns
add_content_slide(prs, "Cross-Department Patterns", [
    "## Recurring Themes Across All 10 Interviews",
    "",
    "Discord overload → CEO, Social (100 msgs), CSO (6 platforms), KOL",
    "Research time sink → Social, PR (10-20 hrs), Paid Growth",
    "Manual reporting → AMs, PR (placement tracking), CEO visibility",
    "Notion adoption friction → Copywriting, AMs, Visual Production",
    "Tool fragmentation → All client-facing roles",
    "",
    "## Key Insight",
    "Solutions addressing Discord + Notion + Reporting have compounding impact across 8+ departments"
])

# SLIDE 12: Automation Tiers
add_content_slide(prs, "Automation Opportunities - Tiered Approach", [
    "## Tier 1: Quick Wins (Week 1-2)",
    "Discord Daily Digest, Notion Forms, Task Notifications",
    "",
    "## Tier 2: Core Automation (Week 3-6)",
    "Report Automation, Research Briefs, Calendar Sync",
    "",
    "## Tier 3: Intelligence Layer (Month 2-3)",
    "MCP Servers, Reply Finder, Sentiment Analysis"
])

# SLIDE 13: Tier 1 Detail
add_content_slide(prs, "Tier 1: Quick Wins (Week 1-2)", [
    "## 1. Discord Daily Digest",
    "AI summarizes key messages by channel",
    "Highlights action items and decisions",
    "Impact: 70% reduction in CEO Discord time",
    "",
    "## 2. Notion Forms with Required Fields",
    "Mandatory fields for creative briefs",
    "Impact: Eliminate 2+ hour delays on tickets",
    "",
    "## 3. Auto-Create Tasks on Assignment",
    "Personal calendar sync when assigned",
    "Impact: Reduce 'did you see my message?' pings"
])

# SLIDE 14: Tier 2 Detail
add_content_slide(prs, "Tier 2: Core Automation (Week 3-6)", [
    "## 4. Automated W/M/Q Reports",
    "Pull data from content calendars",
    "AI generates performance summaries",
    "Impact: Save Diana 2-4 hrs/week, scale to all AMs",
    "",
    "## 5. Research Brief Generator",
    "Morning briefs per client niche",
    "Aggregate crypto news by topic",
    "Impact: 10-min scroll = daily research done",
    "",
    "## 6. Content Calendar Automation",
    "'Done' status → auto-sync to calendar"
])

# SLIDE 15: Tier 3 Detail
add_content_slide(prs, "Tier 3: Intelligence Layer (Month 2-3)", [
    "## 7. Reply Opportunity Finder",
    "Monitor relevant conversations per client",
    "AI suggests 5 reply angles (not full replies)",
    "Impact: Solve 'biggest time sink' for social team",
    "",
    "## 8. VP Performance Feedback Loop",
    "Notify creators when content performs well",
    "Impact: Motivation + learning loop",
    "",
    "## 9. Client Communication Intelligence",
    "Sentiment analysis across Slack/Telegram/Discord",
    "Impact: Proactive client management"
])

# SLIDE 16: Technical Architecture
add_content_slide(prs, "Technical Architecture", [
    "## AI Layer",
    "Claude / GPT for intelligence",
    "",
    "## MCP Servers (Custom Built)",
    "Discord Intel | Notion Ops | Research Hub",
    "",
    "## n8n Workflow Engine",
    "Orchestration & Automation",
    "",
    "## Integrations",
    "Discord | Notion | X/Twitter APIs | Google Workspace"
])

# SLIDE 17: MCP Components
add_content_slide(prs, "MCP Server Components", [
    "## 1. Discord Intelligence MCP",
    "Channel summarization, action items, priority filtering",
    "",
    "## 2. Notion Operations MCP",
    "Required field validation, stage triggers, report aggregation",
    "",
    "## 3. Research Hub MCP",
    "Multi-source news, client profile matching, trend detection"
])

# SLIDE 18: Implementation Roadmap
add_content_slide(prs, "12-Week Implementation Roadmap", [
    "## Phase 1: Foundation (Weeks 1-2)",
    "Discord digest, Notion forms, n8n setup",
    "",
    "## Phase 2: Core Automation (Weeks 3-6)",
    "Report automation, research briefs, calendar sync",
    "",
    "## Phase 3: Intelligence (Weeks 7-10)",
    "MCP servers, reply finder, sentiment analysis",
    "",
    "## Phase 4: Optimization (Weeks 11-12)",
    "Training, refinement, documentation"
])

# SLIDE 19: Phase 1 Detail
add_content_slide(prs, "Phase 1: Foundation (Weeks 1-2)", [
    "## Deliverables",
    "n8n instance configured",
    "Discord bot deployed",
    "Daily digest workflow live",
    "Notion forms with validation",
    "",
    "## Success Metrics",
    "CEO Discord time reduced by 50%+",
    "Zero incomplete creative briefs",
    "Team adoption > 80%",
    "",
    "## Investment: $5,000 - $8,000"
])

# SLIDE 20: Phase 2 Detail
add_content_slide(prs, "Phase 2: Core Automation (Weeks 3-6)", [
    "## Deliverables",
    "W/M/Q report automation",
    "Research brief generator",
    "Content calendar sync",
    "AM dashboard",
    "",
    "## Success Metrics",
    "Report generation: 4 hrs → 15 mins",
    "Research time: 50% reduction",
    "Calendar adoption > 90%",
    "",
    "## Investment: $12,000 - $18,000"
])

# SLIDE 21: Phase 3 Detail
add_content_slide(prs, "Phase 3: Intelligence Layer (Weeks 7-10)", [
    "## Deliverables",
    "Discord MCP server",
    "Notion MCP server",
    "Research Hub MCP",
    "Reply opportunity finder",
    "",
    "## Success Metrics",
    "AI-assisted responses in < 2 min",
    "Proactive client issue detection",
    "Social team efficiency +40%",
    "",
    "## Investment: $15,000 - $25,000"
])

# SLIDE 22: ROI Projections
add_content_slide(prs, "ROI Projections - Conservative Estimates", [
    "## Weekly Time Savings → Annual Value (@$50/hr)",
    "",
    "Discord Digest: 10 hrs/week → $26,000/year",
    "Report Automation: 8 hrs/week → $20,800/year",
    "Research Briefs: 15 hrs/week → $39,000/year",
    "Brief Validation: 5 hrs/week → $13,000/year",
    "Calendar Sync: 3 hrs/week → $7,800/year",
    "",
    "## TOTAL: 41 hrs/week → $106,600/year",
    "",
    "## Payback Period: 3-4 months"
])

# SLIDE 23: Scaling Impact
add_content_slide(prs, "The Real Value: Enabling Growth", [
    "## Without Automation",
    "$3M → $6M requires 40 → 80+ people",
    "Discord becomes unusable at scale",
    "Management overhead explodes",
    "",
    "## With Automation",
    "$3M → $6M with 40 → 60 people",
    "20 fewer hires needed",
    "At $60K avg salary = $1.2M annual savings",
    "",
    "## Total First-Year Value: $1.3M+"
])

# SLIDE 24: Risk Mitigation
add_two_column_slide(prs, "Risk Mitigation",
    [
        "## Team adoption",
        "Phased rollout, champions per dept",
        "",
        "## Tool fatigue",
        "Integrate into existing tools",
        "",
        "## AI accuracy",
        "Human-in-loop for critical decisions"
    ],
    [
        "## Data security",
        "Self-hosted n8n, no data leaves systems",
        "",
        "## Vendor lock-in",
        "Open-source stack, portable workflows",
        "",
        "## Complexity",
        "Start simple, iterate based on feedback"
    ]
)

# SLIDE 25: Engagement Options
add_content_slide(prs, "Engagement Options", [
    "## Option A: Foundation Only",
    "Weeks 1-2 deliverables",
    "Discord digest + Notion forms",
    "Investment: $5,000 - $8,000",
    "",
    "## Option B: Core Package (Recommended)",
    "Phases 1 + 2",
    "Full report + research automation",
    "Investment: $17,000 - $26,000",
    "",
    "## Option C: Complete Transformation",
    "All 3 phases + MCP ecosystem",
    "Investment: $32,000 - $51,000"
])

# SLIDE 26: Next Steps
add_content_slide(prs, "Next Steps", [
    "## Today",
    "Select engagement option",
    "",
    "## This Week",
    "Kick-off meeting, access provisioning",
    "",
    "## Week 1",
    "Foundation deployment begins",
    "",
    "## Week 2",
    "First results visible (Discord digest live)",
    "",
    "## Ongoing",
    "Bi-weekly progress reviews"
])

# SLIDE 27: Q&A
add_title_slide(prs, "Questions & Discussion", "Priority adjustments?\nSpecific integrations?\nTeam members for pilot?\nTimeline constraints?")

# SLIDE 28: Contact
add_title_slide(prs, "Ready to Scale Impact3 Intelligently", "Reprise AI\nkeshav@reprisesai.com")

# Save
output_path = os.path.join(os.path.expanduser("~"), "Documents", "Impact3_AI_Audit_Presentation.pptx")

try:
    prs.save(output_path)
    print(f"✓ Presentation saved successfully!")
    print(f"  Location: {output_path}")
except Exception as e:
    print(f"✗ Error creating presentation: {e}")
    exit(1)
